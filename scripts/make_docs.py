from pathlib import Path
import re
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Preformatted
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / 'docs'

styles = getSampleStyleSheet()
styles.add(ParagraphStyle(name='CodeBlock', parent=styles['Code'], fontName='Courier', fontSize=7, leading=9, backColor=colors.whitesmoke, borderColor=colors.lightgrey, borderWidth=0.25, borderPadding=4))
styles.add(ParagraphStyle(name='SmallBody', parent=styles['BodyText'], fontSize=9, leading=12, spaceAfter=6))
styles.add(ParagraphStyle(name='DocTitle', parent=styles['Title'], fontSize=20, leading=24, spaceAfter=14))
styles.add(ParagraphStyle(name='DocH1', parent=styles['Heading1'], fontSize=15, leading=18, spaceBefore=10, spaceAfter=6))
styles.add(ParagraphStyle(name='DocH2', parent=styles['Heading2'], fontSize=12, leading=15, spaceBefore=8, spaceAfter=4))

def clean_inline(text: str) -> str:
    text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
    text = re.sub(r'`([^`]+)`', r'<font name="Courier">\1</font>', text)
    text = re.sub(r'\*\*([^*]+)\*\*', r'<b>\1</b>', text)
    return text

def markdown_to_story(md_text: str):
    story = []
    in_code = False
    code_lines = []
    first_heading = True
    for raw in md_text.splitlines():
        line = raw.rstrip()
        if line.startswith('```'):
            if not in_code:
                in_code = True
                code_lines = []
            else:
                in_code = False
                code = '\n'.join(code_lines[:60])
                story.append(Preformatted(code, styles['CodeBlock']))
                story.append(Spacer(1, 0.08 * inch))
            continue
        if in_code:
            code_lines.append(line)
            continue
        if not line:
            story.append(Spacer(1, 0.04 * inch))
            continue
        if line.startswith('# '):
            if not first_heading:
                story.append(PageBreak())
            first_heading = False
            story.append(Paragraph(clean_inline(line[2:]), styles['DocTitle']))
        elif line.startswith('## '):
            story.append(Paragraph(clean_inline(line[3:]), styles['DocH1']))
        elif line.startswith('### '):
            story.append(Paragraph(clean_inline(line[4:]), styles['DocH2']))
        elif line.startswith('|'):
            story.append(Paragraph(clean_inline(line), styles['CodeBlock']))
        elif line.startswith('- ') or line.startswith('1. ') or line.startswith('2. ') or line.startswith('3. ') or line.startswith('4. ') or line.startswith('5. '):
            story.append(Paragraph('&bull; ' + clean_inline(line[2:] if line.startswith('- ') else line), styles['SmallBody']))
        else:
            story.append(Paragraph(clean_inline(line), styles['SmallBody']))
    return story

def build_pdf(md_name: str, pdf_name: str):
    md_text = (DOCS / md_name).read_text()
    doc = SimpleDocTemplate(str(DOCS / pdf_name), pagesize=letter, rightMargin=0.65*inch, leftMargin=0.65*inch, topMargin=0.65*inch, bottomMargin=0.65*inch)
    story = markdown_to_story(md_text)
    doc.build(story)

build_pdf('ARCHITECTURE.md', 'ARCHITECTURE.pdf')
build_pdf('SECURITY_AUDIT.md', 'SECURITY_AUDIT.pdf')
build_pdf('GAS_REPORT.md', 'GAS_REPORT.pdf')

# Build final presentation deck.
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(0.65)
    top = Inches(0.35)
    title_box = slide.shapes.add_textbox(left, top, Inches(12), Inches(0.7))
    tf = title_box.text_frame
    tf.text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True
    body = slide.shapes.add_textbox(left, Inches(1.35), Inches(12), Inches(5.6))
    bf = body.text_frame
    bf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(21)
        p.space_after = Pt(10)
    return slide

slides = [
    ('GameFi Economy - Option B', ['Full-stack decentralized protocol for game items, resources, rentals, loot, and DAO governance.', 'Built with Foundry, React, The Graph, Chainlink-compatible mocks, and L2 deployment scripts.', 'Team split: game mechanics, economy/governance, integration/delivery.']),
    ('Problem and user flows', ['Players need resource swaps, crafting, loot boxes, and item rentals.', 'Liquidity providers earn exposure through LP shares.', 'DAO voters govern recipes, drop rates, treasury, and upgrades.']),
    ('System architecture', ['ERC-1155 item economy: GameItems1155.', 'ERC-20 resources: GOLD and CRYSTAL.', 'AMM, rental vault, ERC-4626 treasury, UUPS parameter proxy.', 'The Graph indexes events; frontend reads indexed data and contracts.']),
    ('Game mechanics', ['Crafting burns resource tokens and mints ERC-1155 items.', 'LootDrop burns loot boxes and uses VRF-compatible randomness.', 'RentalVault escrows items and uses pull-payment earnings.']),
    ('Economy layer', ['AMMPool is a from-scratch x*y=k pool with 0.3% fee.', 'LP token is built into the AMM contract.', 'GameVault4626 acts as Timelock-controlled treasury vault.', 'PriceFeedAdapter rejects stale or invalid Chainlink prices.']),
    ('Governance', ['GameToken is ERC20Votes + ERC20Permit.', 'Governor settings: 1-day delay, 1-week period, 4% quorum, 1% proposal threshold.', 'Timelock delay: 2 days.', 'Lifecycle test demonstrates propose -> vote -> queue -> execute.']),
    ('Advanced Solidity', ['UUPS upgrade: GameParametersV1 to GameParametersV2.', 'Factory uses CREATE and CREATE2.', 'AMMMath includes Solidity and inline Yul quote paths for benchmarking.']),
    ('Security design', ['AccessControl on privileged functions.', 'ReentrancyGuard and CEI on external transfer paths.', 'SafeERC20 for ERC-20 movements.', 'No tx.origin, no transfer/send, no block values as randomness.']),
    ('Testing strategy', ['88 test functions plus 6 invariants.', 'Unit, fuzz, invariant, fork, and vulnerability case-study tests.', 'Fuzz coverage includes AMM swaps, vault operations, and voting power.', 'Invariants include k, reserves, LP supply, treasury accounting, and vault rounding.']),
    ('Frontend and subgraph', ['Frontend supports wallet connection, network detection, readable errors, balances, votes, reserves, vault shares.', 'State-changing actions: swap, craft, deposit, delegate, vote.', 'Subgraph indexes players, crafts, items, swaps, rentals, proposals, and votes.']),
    ('Deployment and demo', ['Deploy with script/Deploy.s.sol to Base Sepolia, Arbitrum Sepolia, or Optimism Sepolia.', 'Run VerifyPostDeploy to check Timelock and Governor settings.', 'Fill deployment JSON, frontend env, subgraph addresses, gas table, and final explorer links.'])
]
for title, bullets in slides:
    add_slide(title, bullets)

pptx_path = DOCS / 'Final_Presentation.pptx'
prs.save(pptx_path)
print('[OK] wrote docs PDFs and PPTX')
